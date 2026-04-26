from __future__ import annotations

import io
import os
import shutil
import subprocess
import tempfile
from pathlib import Path
from typing import Any

from src.config import settings
from src.utils.logging import get_logger

logger = get_logger(__name__)


def parse_document(file_path: str) -> dict[str, Any]:
    """
    Parse documents: PDF, Word (.docx/.doc), PowerPoint (.pptx/.ppt),
    Excel (.xlsx/.xls), and plain text.

    Legacy binary formats (.doc, .ppt, .xls) are converted to their modern
    XML equivalents via LibreOffice when available, then parsed normally.
    """
    file_ext = Path(file_path).suffix.lower()

    try:
        if file_ext == ".txt":
            return _parse_text(file_path)
        elif file_ext == ".pdf":
            return _parse_pdf(file_path)
        elif file_ext == ".docx":
            return _parse_docx(file_path)
        elif file_ext == ".doc":
            return _parse_legacy(file_path, target_ext=".docx", parser=_parse_docx)
        elif file_ext == ".pptx":
            return _parse_pptx(file_path)
        elif file_ext == ".ppt":
            return _parse_legacy(file_path, target_ext=".pptx", parser=_parse_pptx)
        elif file_ext == ".xlsx":
            return _parse_xlsx(file_path)
        elif file_ext == ".xls":
            return _parse_legacy(file_path, target_ext=".xlsx", parser=_parse_xlsx)
        else:
            raise ValueError(f"Unsupported file format: {file_ext}")
    except Exception as e:
        logger.error(f"Failed to parse document {file_path}: {e}")
        raise


# ── Legacy format conversion ──────────────────────────────────────────────────

def _parse_legacy(file_path: str, target_ext: str, parser) -> dict[str, Any]:
    """
    Convert a legacy binary format (.doc/.ppt/.xls) to its modern equivalent
    using LibreOffice headless, then parse with the appropriate parser.
    Raises ValueError if LibreOffice is not installed.
    """
    libreoffice = shutil.which("libreoffice") or shutil.which("soffice")
    if not libreoffice:
        raise ValueError(
            f"Legacy format {Path(file_path).suffix} requires LibreOffice to convert. "
            f"Install it with: brew install libreoffice  "
            f"Or save the file as {target_ext} and re-upload."
        )

    with tempfile.TemporaryDirectory() as tmp_dir:
        result = subprocess.run(
            [libreoffice, "--headless", "--convert-to",
             target_ext.lstrip("."), "--outdir", tmp_dir, file_path],
            capture_output=True, text=True, timeout=120,
        )
        if result.returncode != 0:
            raise RuntimeError(
                f"LibreOffice conversion failed: {result.stderr.strip()}"
            )

        converted = Path(tmp_dir) / (Path(file_path).stem + target_ext)
        if not converted.exists():
            raise RuntimeError(f"Converted file not found: {converted}")

        logger.info(f"Converted {file_path} → {converted}")
        parsed = parser(str(converted))
        # Restore original filename in result
        parsed["filename"] = os.path.basename(file_path)
        return parsed


# ── Text ──────────────────────────────────────────────────────────────────────

def _parse_text(file_path: str) -> dict[str, Any]:
    with open(file_path, "r", encoding="utf-8") as f:
        content = f.read()
    logger.info(f"Parsed text file: {file_path}")
    return {"filename": os.path.basename(file_path), "content": content, "format": "text", "images": []}


# ── PDF ───────────────────────────────────────────────────────────────────────

def _parse_pdf(file_path: str) -> dict[str, Any]:
    """Extract text and embedded images from PDF using pypdf."""
    from pypdf import PdfReader

    content_parts: list[str] = []
    images: list[str] = []
    base = os.path.splitext(os.path.basename(file_path))[0]

    with open(file_path, "rb") as f:
        reader = PdfReader(f)
        for page_num, page in enumerate(reader.pages):
            text = page.extract_text()
            if text:
                content_parts.append(f"--- Page {page_num + 1} ---\n{text}")
            try:
                for img_obj in page.images:
                    path = _save_image(img_obj.data, base, len(images))
                    if path:
                        images.append(path)
            except Exception as e:
                logger.warning(f"PDF image extraction page {page_num + 1}: {e}")

    content = "\n\n".join(content_parts)
    logger.info(f"Parsed PDF: {file_path} ({len(reader.pages)} pages, {len(images)} images)")
    return {
        "filename": os.path.basename(file_path),
        "content": content,
        "format": "pdf",
        "pages": len(reader.pages),
        "images": images,
    }


# ── Word ──────────────────────────────────────────────────────────────────────

def _parse_docx(file_path: str) -> dict[str, Any]:
    """Extract text, tables, and embedded images from .docx."""
    from docx import Document

    doc = Document(file_path)
    content_parts: list[str] = []
    images: list[str] = []
    base = os.path.splitext(os.path.basename(file_path))[0]

    for para in doc.paragraphs:
        if para.text.strip():
            content_parts.append(para.text)

    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text for cell in row.cells if cell.text.strip())
            if row_text:
                content_parts.append(row_text)

    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            try:
                path = _save_image(rel.target_part.blob, base, len(images))
                if path:
                    images.append(path)
            except Exception as e:
                logger.warning(f"Word image extraction: {e}")

    content = "\n\n".join(content_parts)
    logger.info(f"Parsed Word: {file_path} ({len(images)} images)")
    return {
        "filename": os.path.basename(file_path),
        "content": content,
        "format": "word",
        "images": images,
    }


# ── PowerPoint ────────────────────────────────────────────────────────────────

def _parse_pptx(file_path: str) -> dict[str, Any]:
    """Extract text and images from each slide in a .pptx file."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(file_path)
    content_parts: list[str] = []
    images: list[str] = []
    base = os.path.splitext(os.path.basename(file_path))[0]

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_content = [f"--- Slide {slide_num} ---"]

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_content.append(shape.text)

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    path = _save_image(shape.image.blob, f"{base}_s{slide_num}", len(images))
                    if path:
                        images.append(path)
                except Exception as e:
                    logger.warning(f"PPT image slide {slide_num}: {e}")

        if len(slide_content) > 1:
            content_parts.append("\n".join(slide_content))

    content = "\n\n".join(content_parts)
    logger.info(f"Parsed PowerPoint: {file_path} ({len(prs.slides)} slides, {len(images)} images)")
    return {
        "filename": os.path.basename(file_path),
        "content": content,
        "format": "powerpoint",
        "slides": len(prs.slides),
        "images": images,
    }


# ── Excel ─────────────────────────────────────────────────────────────────────

def _parse_xlsx(file_path: str) -> dict[str, Any]:
    """Extract all cell values from all sheets in .xlsx."""
    from openpyxl import load_workbook

    wb = load_workbook(file_path, data_only=True)
    content_parts: list[str] = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        sheet_content = [f"--- Sheet: {sheet_name} ---"]
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None and str(c).strip()]
            if cells:
                sheet_content.append(" | ".join(cells))
        if len(sheet_content) > 1:
            content_parts.append("\n".join(sheet_content))

    content = "\n\n".join(content_parts)
    logger.info(f"Parsed Excel: {file_path} ({len(wb.sheetnames)} sheets)")
    return {
        "filename": os.path.basename(file_path),
        "content": content,
        "format": "excel",
        "sheets": len(wb.sheetnames),
        "images": [],
    }


# ── Helpers ───────────────────────────────────────────────────────────────────

def _save_image(image_bytes: bytes, base_name: str, index: int) -> str | None:
    """Save raw image bytes as PNG to a temp file. Returns path or None."""
    try:
        from PIL import Image
        img = Image.open(io.BytesIO(image_bytes))
        path = os.path.join(tempfile.gettempdir(), f"{base_name}_img{index}.png")
        img.save(path, "PNG")
        return path
    except Exception as e:
        logger.warning(f"Could not save image {index}: {e}")
        return None


def extract_text_from_image(image_path: str) -> str:
    """
    Extract text from an image using configured method: OCR (pytesseract) or VLM.

    Returns extracted text content that can be appended to document text.
    """
    if settings.IMAGE_EXTRACTION_MODE == "vlm":
        # VLM mode: extract structured entities and relationships
        result = _extract_with_vlm(image_path)
        text_parts = []

        for entity in result.get("entities", []):
            desc = entity.get("description", "")
            text_parts.append(f"Entity: {entity.get('name')} ({entity.get('type')}) - {desc}")

        for rel in result.get("relationships", []):
            desc = rel.get("description", "")
            text_parts.append(
                f"Relationship: {rel.get('from')} --[{rel.get('type')}]--> {rel.get('to')} - {desc}"
            )

        if text_parts:
            return "\n".join(text_parts)
        elif result.get("note"):
            logger.info(f"VLM: {result.get('note')} for {image_path}")
            return ""
        elif result.get("error") or result.get("parse_error"):
            logger.warning(f"VLM extraction failed for {image_path}, falling back to OCR")
            return _extract_with_ocr(image_path)
        return ""
    else:
        # Default: OCR mode (pytesseract)
        return _extract_with_ocr(image_path)


def _extract_with_ocr(image_path: str) -> str:
    """Extract text from image using pytesseract OCR."""
    try:
        import pytesseract
        from PIL import Image as PILImage
        ocr_text = pytesseract.image_to_string(PILImage.open(image_path)).strip()
        if ocr_text:
            logger.info(f"[OCR] Extracted {len(ocr_text)} chars from {image_path}")
            return ocr_text
        return ""
    except Exception as e:
        logger.warning(f"[OCR] Failed for {image_path}: {e}")
        return ""


def _extract_with_vlm(image_path: str) -> dict[str, Any]:
    """Extract entities/relationships from image using Azure OpenAI vision."""
    try:
        from src.ingestion.vlm_extractor import extract_entities_from_image
        return extract_entities_from_image(image_path)
    except Exception as e:
        logger.warning(f"[VLM] Failed for {image_path}: {e}")
        return {"entities": [], "relationships": [], "error": str(e)}
